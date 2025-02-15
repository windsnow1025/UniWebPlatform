import logging
from io import BytesIO

import docx
import fitz
import httpx
import openpyxl
from docx.oxml import CT_P, CT_Tbl
from fastapi import HTTPException
from pptx import Presentation

from chat.logic.message_preprocess import file_type_checker


async def extract_text_from_file(file_url: str) -> str:
    async with httpx.AsyncClient() as client:
        try:
            response = await client.get(file_url)
        except Exception as e:
            logging.exception(e)
            raise HTTPException(status_code=500, detail=str(e))

        if response.status_code != 200:
            status_code = response.status_code
            text = response.text
            raise HTTPException(status_code=status_code, detail=text)

        file_content = response.content

        if file_type_checker.is_file_type_supported(file_url) is False:
            raise HTTPException(status_code=415, detail="legacy filetypes ('.doc', '.xls', '.ppt') are not supported")
        file_type, sub_type = file_type_checker.get_file_type(file_url)
        if sub_type == "code":
            return extract_text_from_code(file_content)
        if sub_type == "pdf":
            return extract_text_from_pdf(file_content)
        if sub_type == "word":
            return extract_text_from_word(file_content)
        if sub_type == "excel":
            return extract_text_from_excel(file_content)
        if sub_type == "ppt":
            return extract_text_from_ppt(file_content)

        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError as e:
            logging.exception(e)
            raise HTTPException(status_code=415, detail=str(e))


def extract_text_from_code(file_content: bytes) -> str:
    return file_content.decode('utf-8')


def extract_text_from_pdf(file_content: bytes) -> str:
    text = ""
    with fitz.open(stream=file_content, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text


def extract_text_from_word(file_content: bytes) -> str:
    text = ""
    try:
        doc = docx.Document(BytesIO(file_content))

        # Extract text from headers
        for section in doc.sections:
            header = section.header
            for para in header.paragraphs:
                text += para.text + "\n"

        # Extract text from the main document body
        for element in doc.element.body:
            if isinstance(element, CT_P):
                para = docx.text.paragraph.Paragraph(element, doc)
                text += para.text + "\n"
            elif isinstance(element, CT_Tbl):
                table = docx.table.Table(element, doc)
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"

        # Extract text from footers
        for section in doc.sections:
            footer = section.footer
            for para in footer.paragraphs:
                text += para.text + "\n"

        return text
    except Exception as e:
        logging.exception(e)
        raise HTTPException(status_code=415, detail=str(e))


def extract_text_from_excel(file_content: bytes) -> str:
    text = ""
    try:
        workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
        for sheet in workbook:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
        return text
    except Exception as e:
        logging.exception(e)
        raise HTTPException(status_code=415, detail=str(e))


def extract_text_from_ppt(file_content: bytes) -> str:
    text = ""
    try:
        presentation = Presentation(BytesIO(file_content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logging.exception(e)
        raise HTTPException(status_code=415, detail=str(e))
